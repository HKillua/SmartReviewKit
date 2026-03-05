"""PPTX (PowerPoint) document loader with OMML formula extraction."""

from __future__ import annotations

import hashlib
import logging
import re
from pathlib import Path
from typing import Any

from src.core.types import Document, DocumentSection
from src.libs.loader.base_loader import BaseLoader

logger = logging.getLogger(__name__)

try:
    from lxml import etree
    LXML_AVAILABLE = True
except ImportError:
    LXML_AVAILABLE = False

_OMML_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _omml_tag(local: str) -> str:
    return f"{{{_OMML_NS}}}{local}"


_EXERCISE_RE = re.compile(r"例题|练习|习题|Exercise|Problem|思考题", re.IGNORECASE)
_DEFINITION_RE = re.compile(r"定义|Definition", re.IGNORECASE)
_THEOREM_RE = re.compile(r"定理|Theorem|引理|Lemma", re.IGNORECASE)
_EXAMPLE_RE = re.compile(r"[例示]例|Example|示例", re.IGNORECASE)


def _infer_content_type(text: str) -> str:
    if _EXERCISE_RE.search(text[:200]):
        return "exercise"
    if _DEFINITION_RE.search(text[:200]):
        return "definition"
    if _THEOREM_RE.search(text[:200]):
        return "theorem"
    if _EXAMPLE_RE.search(text[:200]):
        return "example"
    return "concept"


class PptxLoader(BaseLoader):
    """Load .pptx files with OMML formula extraction and structured sections."""

    def __init__(self, extract_images: bool = False, image_storage_dir: str = "data/images") -> None:
        self._extract_images = extract_images
        self._image_dir = Path(image_storage_dir)

    def load(self, file_path: str | Path) -> Document:
        path = self._validate_file(file_path)

        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError("python-pptx is required: pip install python-pptx") from exc

        prs = Presentation(str(path))
        doc_hash = hashlib.sha256(path.read_bytes()).hexdigest()[:12]

        parts: list[str] = []
        images: list[dict[str, Any]] = []
        sections: list[dict[str, Any]] = []
        image_seq = 0

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_title = ""
            slide_texts: list[str] = []
            slide_images: list[dict[str, Any]] = []

            for shape in slide.shapes:
                # Text with interleaved OMML formulas
                if shape.has_text_frame:
                    is_title = (
                        slide.shapes.title is not None
                        and shape.shape_id == slide.shapes.title.shape_id
                    )
                    para_texts = self._extract_text_with_math(shape)
                    combined = "\n".join(para_texts).strip()
                    if not combined:
                        continue
                    if is_title:
                        slide_title = combined
                    else:
                        slide_texts.append(combined)

                if shape.has_table:
                    slide_texts.append(self._table_to_markdown(shape.table))

                if self._extract_images and shape.shape_type == 13:
                    image_seq += 1
                    img_id = f"{doc_hash}_{slide_idx}_{image_seq}"
                    slide_texts.append(f"[IMAGE: {img_id}]")
                    img_path = self._save_image(shape, img_id)
                    img_meta = {"id": img_id, "page": slide_idx, "path": img_path}
                    images.append(img_meta)
                    slide_images.append(img_meta)

            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            header = f"## Slide {slide_idx}"
            if slide_title:
                header += f": {slide_title}"
            body_lines = [header, ""]
            body_lines.extend(slide_texts)
            if notes_text:
                body_lines.append(f"\n> **备注**: {notes_text}")
            section_text = "\n".join(body_lines)
            parts.append(section_text)

            content_body = "\n".join(slide_texts)
            has_formula = "$" in content_body
            ct = _infer_content_type(slide_title + " " + content_body)
            if has_formula and ct == "concept" and len(content_body) < 300:
                ct = "formula"

            sections.append(DocumentSection(
                title=slide_title,
                level=1 if slide_idx == 1 and not slide_texts else 2,
                content=content_body,
                content_type=ct,
                page_or_slide=slide_idx,
                has_formula=has_formula,
                images=slide_images,
            ).to_dict())

        full_text = f"# {path.stem}\n\n" + "\n\n---\n\n".join(parts)

        return Document(
            id=doc_hash,
            text=full_text,
            metadata={
                "source_path": str(path),
                "doc_type": "pptx",
                "title": path.stem,
                "slide_count": len(prs.slides),
                "images": images,
                "sections": sections,
            },
        )

    # ------------------------------------------------------------------
    # Text + OMML extraction
    # ------------------------------------------------------------------

    def _extract_text_with_math(self, shape: Any) -> list[str]:
        """Extract paragraphs from a shape, interleaving OMML formulas as LaTeX."""
        paragraphs: list[str] = []
        for para in shape.text_frame.paragraphs:
            parts = self._parse_paragraph_xml(para._element)
            line = "".join(parts).strip()
            if line:
                paragraphs.append(line)
        return paragraphs

    def _parse_paragraph_xml(self, p_elem: Any) -> list[str]:
        """Walk <a:p> children for <a:r> (text runs) and <m:oMath> (formulas)."""
        parts: list[str] = []
        if not LXML_AVAILABLE:
            if hasattr(p_elem, "text") and p_elem.text:
                return [p_elem.text]
            return parts

        for child in p_elem:
            tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ""
            ns = etree.QName(child.tag).namespace if isinstance(child.tag, str) else ""

            if ns == _OMML_NS and tag in ("oMath", "oMathPara"):
                from src.libs.loader.math_utils import omml_to_latex
                latex = omml_to_latex(child).strip()
                if latex:
                    if tag == "oMathPara":
                        parts.append(f"$${latex}$$")
                    else:
                        parts.append(f"${latex}$")
            elif ns == _A_NS and tag == "r":
                t_elem = child.find(f"{{{_A_NS}}}t")
                if t_elem is not None and t_elem.text:
                    parts.append(t_elem.text)

        return parts

    # ------------------------------------------------------------------
    # Image saving
    # ------------------------------------------------------------------

    def _save_image(self, shape: Any, img_id: str) -> str:
        self._image_dir.mkdir(parents=True, exist_ok=True)
        blob = shape.image.blob
        ext = shape.image.content_type.split("/")[-1] if shape.image.content_type else "png"
        if ext == "jpeg":
            ext = "jpg"
        out_path = self._image_dir / f"{img_id}.{ext}"
        out_path.write_bytes(blob)
        return str(out_path)

    # ------------------------------------------------------------------
    # Table to Markdown
    # ------------------------------------------------------------------

    @staticmethod
    def _table_to_markdown(table: Any) -> str:
        rows: list[list[str]] = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        if not rows:
            return ""
        lines = ["| " + " | ".join(rows[0]) + " |"]
        lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
        for row in rows[1:]:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines)
